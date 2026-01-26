"""
Extract text content from PDF and PPTX files for review
"""
import os

def extract_pdf_text(pdf_path):
    """Extract text from PDF file"""
    try:
        import PyPDF2
        text = []
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            print(f"PDF has {num_pages} pages")
            
            for i, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                text.append(f"\n\n========== PAGE {i+1} ==========\n\n")
                text.append(page_text)
        
        return ''.join(text)
    except ImportError:
        return "ERROR: PyPDF2 not installed. Install with: pip install PyPDF2"
    except Exception as e:
        return f"ERROR extracting PDF: {str(e)}"

def extract_pptx_text(pptx_path):
    """Extract text from PowerPoint file"""
    try:
        from pptx import Presentation
        text = []
        prs = Presentation(pptx_path)
        
        print(f"PowerPoint has {len(prs.slides)} slides")
        
        for i, slide in enumerate(prs.slides):
            text.append(f"\n\n========== SLIDE {i+1} ==========\n\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text + "\n")
        
        return ''.join(text)
    except ImportError:
        return "ERROR: python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"ERROR extracting PPTX: {str(e)}"

if __name__ == "__main__":
    base_dir = r"c:\Users\msibr\Documents\MCA\SEM 4\Project\kabaddi_trainer"
    
    # Extract PDF
    pdf_path = os.path.join(base_dir, "documents", "Report.pdf")
    print(f"Extracting PDF: {pdf_path}")
    pdf_text = extract_pdf_text(pdf_path)
    
    pdf_output = os.path.join(base_dir, "documents", "Report_extracted.txt")
    with open(pdf_output, 'w', encoding='utf-8') as f:
        f.write(pdf_text)
    print(f"PDF text saved to: {pdf_output}\n")
    
    # Extract PPTX
    pptx_path = os.path.join(base_dir, "review1", "2024178011 review-1 [Autosaved].pptx")
    print(f"Extracting PPTX: {pptx_path}")
    pptx_text = extract_pptx_text(pptx_path)
    
    pptx_output = os.path.join(base_dir, "review1", "Presentation_extracted.txt")
    with open(pptx_output, 'w', encoding='utf-8') as f:
        f.write(pptx_text)
    print(f"PPTX text saved to: {pptx_output}")
    
    print("\n=== EXTRACTION COMPLETE ===")
